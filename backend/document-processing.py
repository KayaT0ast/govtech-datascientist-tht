"""
document-processing.py

DocumentProcessor — converts a directory of KB documents into overlapping
text chunks and writes kb_chunks.json.

Supported formats:
    .pdf  — pdfplumber (text + pipe-delimited tables)
    .pptx — python-pptx (slide titles as headings + body text + tables)
    .docx — python-docx (paragraphs with heading styles + tables)
    .txt / .md — plain read

Outputs:
    kb_chunks.json  — list of {source, chunk_id, text} objects

Requires:
    pip install pdfplumber python-pptx python-docx
"""

import json
import re
from pathlib import Path

# ── Chunking constants ────────────────────────────────────────────────────────

_CHUNK_TOKEN_TARGET = 500
_CHUNK_OVERLAP_TOKENS = 50
_WORDS_PER_TOKEN = 0.75
_SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md"}


# ── Text chunker ──────────────────────────────────────────────────────────────

def _chunk_text(text: str, source: str) -> list[dict]:
    """Split text into overlapping chunks, preferring heading boundaries."""
    chunks = []
    sections = re.split(r"\n(?=#{1,3} )", text)
    if len(sections) <= 1:
        sections = [text]

    chunk_id = 0
    target_words = int(_CHUNK_TOKEN_TARGET / _WORDS_PER_TOKEN)
    overlap_words = int(_CHUNK_OVERLAP_TOKENS / _WORDS_PER_TOKEN)

    for section in sections:
        words = section.split()
        if not words:
            continue
        if len(words) <= target_words:
            chunks.append({"source": source, "chunk_id": chunk_id, "text": section.strip()})
            chunk_id += 1
        else:
            start = 0
            while start < len(words):
                end = min(start + target_words, len(words))
                chunks.append({"source": source, "chunk_id": chunk_id,
                                "text": " ".join(words[start:end])})
                chunk_id += 1
                start += target_words - overlap_words
    return chunks


# ── Document converters ───────────────────────────────────────────────────────

def _convert_pdf(path: Path) -> str:
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pip install pdfplumber")

    parts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                parts.append(text.strip())
            for table in page.extract_tables():
                if not table:
                    continue
                rows = [
                    " | ".join(str(cell).strip() if cell is not None else "" for cell in row)
                    for row in table
                ]
                if rows:
                    parts.append("\n".join(rows))
    return "\n\n".join(parts)


def _convert_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("pip install python-pptx")

    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text.strip()
        heading = f"## Slide {i}: {title_text}" if title_text else f"## Slide {i}"
        lines = [heading]
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        lines.append(para.text.strip())
            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(cells))
        if len(lines) > 1:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)


def _convert_docx(path: Path) -> str:
    try:
        import docx
    except ImportError:
        raise RuntimeError("pip install python-docx")

    _HEADING_MAP = {
        "Heading 1": "# ", "Heading 2": "## ", "Heading 3": "### ",
        "Title": "# ", "Subtitle": "## ",
    }
    doc = docx.Document(str(path))
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        prefix = _HEADING_MAP.get(para.style.name, "")
        parts.append(f"{prefix}{text}")
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


_CONVERTERS = {
    ".pdf":  _convert_pdf,
    ".pptx": _convert_pptx,
    ".docx": _convert_docx,
}


# ── Processor ─────────────────────────────────────────────────────────────────

class DocumentProcessor:
    """
    Converts all supported documents in a directory into overlapping text
    chunks and writes kb_chunks.json.

    Usage:
        processor = DocumentProcessor()
        result = processor.run("./knowledge-base")
        # result["chunks"] — list of {source, chunk_id, text}
    """

    def run(self, kb_dir: str | Path, output_path: str | Path | None = None) -> dict:
        """
        Process all supported documents in kb_dir.

        Args:
            kb_dir:      directory containing KB documents
            output_path: destination for kb_chunks.json
                         (defaults to <kb_dir>/kb_chunks.json)

        Returns:
            {"chunks": [...], "processed": [...], "skipped": [...], "output_path": str}
        """
        kb_dir = Path(kb_dir)
        output_path = Path(output_path) if output_path else kb_dir / "kb_chunks.json"

        chunks: list[dict] = []
        processed: list[str] = []
        skipped: list[str] = []

        for path in sorted(kb_dir.iterdir()):
            if not path.is_file():
                continue
            ext = path.suffix.lower()
            if ext not in _SUPPORTED_EXTENSIONS:
                continue
            try:
                if ext in (".txt", ".md"):
                    text = path.read_text(encoding="utf-8", errors="replace")
                else:
                    text = _CONVERTERS[ext](path)
                file_chunks = _chunk_text(text, path.name)
                chunks.extend(file_chunks)
                processed.append(path.name)
                print(f"  {path.name}: {len(file_chunks)} chunks")
            except Exception as e:
                skipped.append(path.name)
                print(f"  [skip] {path.name}: {e}")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(
            json.dumps(chunks, indent=2, ensure_ascii=False), encoding="utf-8"
        )
        print(f"  Total: {len(chunks)} chunks from {len(processed)} file(s)")
        return {"chunks": chunks, "processed": processed,
                "skipped": skipped, "output_path": str(output_path)}
